"""Utility functions for creating charts in PowerPoint presentations.

This module provides functionality to create various chart types in PowerPoint slides
using python-pptx's chart capabilities.
"""

import logging
from typing import Dict, Any, Optional

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


# Mapping of chart type strings to python-pptx chart types
CHART_TYPE_MAP = {
    'bar': XL_CHART_TYPE.BAR_CLUSTERED,
    'bar_stacked': XL_CHART_TYPE.BAR_STACKED,
    'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
    'column_stacked': XL_CHART_TYPE.COLUMN_STACKED,
    'line': XL_CHART_TYPE.LINE,
    'line_markers': XL_CHART_TYPE.LINE_MARKERS,
    'pie': XL_CHART_TYPE.PIE,
    'doughnut': XL_CHART_TYPE.DOUGHNUT,
    'area': XL_CHART_TYPE.AREA,
    'area_stacked': XL_CHART_TYPE.AREA_STACKED,
    'scatter': XL_CHART_TYPE.XY_SCATTER,
    'radar': XL_CHART_TYPE.RADAR,
}


class ChartDataError(Exception):
    """Exception raised when chart data is invalid."""
    pass


def validate_chart_data(chart_data: Dict[str, Any], chart_type: str) -> None:
    """Validate chart data structure.

    Args:
        chart_data: Dictionary containing categories and series.
        chart_type: Type of chart being created.

    Raises:
        ChartDataError: If data is invalid.
    """
    if not chart_data:
        raise ChartDataError("Chart data is required")

    if chart_type not in CHART_TYPE_MAP:
        raise ChartDataError(f"Unknown chart type: {chart_type}. Available: {', '.join(CHART_TYPE_MAP.keys())}")

    # Scatter charts don't use categories
    if chart_type != 'scatter':
        if 'categories' not in chart_data:
            raise ChartDataError("Chart data must include 'categories'")
        if not chart_data['categories']:
            raise ChartDataError("Categories list cannot be empty")

    if 'series' not in chart_data:
        raise ChartDataError("Chart data must include 'series'")
    if not chart_data['series']:
        raise ChartDataError("Series list cannot be empty")

    # Validate each series
    for i, series in enumerate(chart_data['series']):
        if not isinstance(series, dict):
            raise ChartDataError(f"Series {i} must be a dictionary")
        if 'name' not in series:
            raise ChartDataError(f"Series {i} must have a 'name'")
        if 'values' not in series:
            raise ChartDataError(f"Series {i} must have 'values'")
        if not series['values']:
            raise ChartDataError(f"Series {i} values cannot be empty")


def create_chart_data(chart_data: Dict[str, Any]) -> CategoryChartData:
    """Create CategoryChartData from structured data.

    Args:
        chart_data: Dictionary with 'categories' and 'series' keys.

    Returns:
        CategoryChartData object ready for chart creation.
    """
    data = CategoryChartData()
    data.categories = chart_data['categories']

    for series in chart_data['series']:
        data.add_series(series['name'], series['values'])

    return data


def add_chart_to_slide(
    slide,
    chart_type: str,
    chart_data: Dict[str, Any],
    left: int,
    top: int,
    width: int,
    height: int,
    has_legend: bool = True,
    legend_position: str = 'right',
    title: Optional[str] = None
) -> None:
    """Add a chart to a slide.

    Args:
        slide: PowerPoint slide object.
        chart_type: Type of chart (bar, line, pie, etc.).
        chart_data: Dictionary with categories and series data.
        left: Left position in EMUs.
        top: Top position in EMUs.
        width: Chart width in EMUs.
        height: Chart height in EMUs.
        has_legend: Whether to show legend.
        legend_position: Legend position (left, right, top, bottom).
        title: Optional chart title.
    """
    logger.debug(f"Adding {chart_type} chart to slide")

    # Validate data
    validate_chart_data(chart_data, chart_type)

    # Get chart type enum
    xl_chart_type = CHART_TYPE_MAP[chart_type]

    # Create chart data
    data = create_chart_data(chart_data)

    # Add chart to slide
    chart_shape = slide.shapes.add_chart(
        xl_chart_type,
        left, top, width, height,
        data
    )

    chart = chart_shape.chart

    # Configure legend
    if has_legend:
        chart.has_legend = True
        legend_positions = {
            'left': XL_LEGEND_POSITION.LEFT,
            'right': XL_LEGEND_POSITION.RIGHT,
            'top': XL_LEGEND_POSITION.TOP,
            'bottom': XL_LEGEND_POSITION.BOTTOM,
        }
        chart.legend.position = legend_positions.get(legend_position, XL_LEGEND_POSITION.RIGHT)
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False

    # Set chart title if provided
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
    else:
        chart.has_title = False

    logger.debug(f"Chart added successfully with {len(chart_data['series'])} series")


# Default chart colors (modern palette)
DEFAULT_CHART_COLORS = [
    RGBColor(0x41, 0x72, 0xC4),  # Blue
    RGBColor(0xED, 0x7D, 0x31),  # Orange
    RGBColor(0xA5, 0xA5, 0xA5),  # Gray
    RGBColor(0xFF, 0xC0, 0x00),  # Yellow
    RGBColor(0x5B, 0x9B, 0xD5),  # Light Blue
    RGBColor(0x70, 0xAD, 0x47),  # Green
    RGBColor(0x9E, 0x48, 0x0E),  # Brown
    RGBColor(0x63, 0x6F, 0x7F),  # Dark Gray
]
