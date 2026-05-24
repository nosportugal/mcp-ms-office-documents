"""PowerPoint helper utilities and slide-building mixin.

This module provides a single SlideHelpers mixin class that consolidates
all common slide operations (text, tables, images) and standalone utility
functions for template loading and data parsing.
"""

import logging
import re
from typing import List, Tuple, Optional, Any

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml

from .constants import (
    BLANK_LAYOUT, CONTENT_LAYOUT,
    DEFAULT_BODY_FONT_SIZE,
    MARGIN_LEFT,
    TABLE_HEADER_FILL, TABLE_HEADER_TEXT, TABLE_ALT_ROW_FILL,
)
from .image_utils import download_image, ImageDownloadError, ImageValidationError
from .inline_formatting import has_inline_formatting, apply_inline_formatting

logger = logging.getLogger(__name__)


# =============================================================================
# Utility Functions
# =============================================================================

# Regex matching a single markdown table separator cell: optional colon, 3+ dashes, optional colon
_SEPARATOR_CELL_RE = re.compile(r'^\s*:?-{3,}:?\s*$')


def _is_separator_row(row: List[str]) -> bool:
    """Check if a row is a markdown table separator row (e.g., |:---|:---:|---:|).

    Uses strict per-cell regex to avoid false positives on content containing dashes.
    """
    return bool(row) and all(_SEPARATOR_CELL_RE.match(cell) for cell in row)


def _extract_alignments(row: List[str]) -> List[Optional[int]]:
    """Extract column alignments from a markdown separator row.

    Args:
        row: List of separator cells (e.g., [':---', ':---:', '---:']).

    Returns:
        List of PP_ALIGN values (LEFT, CENTER, RIGHT) or None per column.
    """
    alignments = []
    for cell in row:
        cell = cell.strip()
        if cell.startswith(':') and cell.endswith(':'):
            alignments.append(PP_ALIGN.CENTER)
        elif cell.endswith(':'):
            alignments.append(PP_ALIGN.RIGHT)
        else:
            alignments.append(None)  # left/default
    return alignments


def parse_table_data(table_data: List[List[str]]) -> tuple:
    """Clean table data by removing markdown separator rows and extracting alignments.

    Args:
        table_data: Raw table data as list of rows.

    Returns:
        Tuple of (cleaned_rows, column_alignments).
        column_alignments is a list of PP_ALIGN values or None per column,
        or None if no separator row was found.
    """
    if not table_data:
        return [], None

    cleaned = []
    col_alignments = None

    for row in table_data:
        if _is_separator_row(row):
            col_alignments = _extract_alignments(row)
        else:
            cleaned.append(row)

    return cleaned, col_alignments


def parse_color(color_hex: str, default: RGBColor) -> RGBColor:
    """Parse hex color string to RGBColor.

    Args:
        color_hex: Hex color string (e.g., "4172C4").
        default: Default color if parsing fails.

    Returns:
        RGBColor object.
    """
    try:
        return RGBColor.from_string(color_hex)
    except (ValueError, AttributeError):
        return default


# =============================================================================
# Consolidated Slide Helpers Mixin
# =============================================================================

class SlideHelpers:
    """Mixin providing all common slide helper methods (text, tables, images).

    Expects the consuming class to have a `self.presentation` attribute
    holding a python-pptx Presentation object.
    """

    # Type hint for IDE — actual attribute is set by the consuming class
    presentation: Any

    # -------------------------------------------------------------------------
    # Slide Management
    # -------------------------------------------------------------------------

    def _get_slide_dimensions(self) -> Tuple[int, int]:
        """Get slide width and height."""
        return self.presentation.slide_width, self.presentation.slide_height

    def _add_blank_slide(self):
        """Add a blank slide and return it."""
        layout = self.presentation.slide_layouts[BLANK_LAYOUT]
        return self.presentation.slides.add_slide(layout)

    def _add_title_content_slide(self, title: str = ""):
        """Add a Title and Content slide and return slide with content placeholder info.

        Args:
            title: Title text for the slide.

        Returns:
            Tuple of (slide, content_left, content_top, content_width, content_height)
        """
        layout = self.presentation.slide_layouts[CONTENT_LAYOUT]
        slide = self.presentation.slides.add_slide(layout)

        # Set title
        if title and len(slide.placeholders) > 0:
            slide.placeholders[0].text = title

        # Get content placeholder bounds (idx 1)
        content_placeholder = None
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 1:
                content_placeholder = placeholder
                break

        if content_placeholder:
            left = content_placeholder.left
            top = content_placeholder.top
            width = content_placeholder.width
            height = content_placeholder.height
            # Remove the placeholder so we can add custom content
            sp = content_placeholder._element
            sp.getparent().remove(sp)
        else:
            # Fallback dimensions
            slide_width, slide_height = self._get_slide_dimensions()
            left = MARGIN_LEFT
            top = Inches(1.5)
            width = slide_width - (2 * MARGIN_LEFT)
            height = slide_height - top - Inches(0.5)

        return slide, left, top, width, height

    def _add_speaker_notes(self, slide, notes_text: Optional[str]) -> None:
        """Add speaker notes to a slide.

        Args:
            slide: PowerPoint slide object.
            notes_text: Text for speaker notes.
        """
        if not notes_text:
            return
        try:
            slide.notes_slide.notes_text_frame.text = notes_text
            logger.debug(f"Added speaker notes: {notes_text[:50]}...")
        except Exception as e:
            logger.warning(f"Could not add speaker notes: {e}")

    # -------------------------------------------------------------------------
    # Text Helpers
    # -------------------------------------------------------------------------

    def _add_text_box(
        self,
        slide,
        text: str,
        left: int,
        top: int,
        width: int,
        height: int,
        font_size: Optional[int] = None,
        bold: bool = False,
        italic: bool = False,
        alignment=PP_ALIGN.LEFT,
        word_wrap: bool = True
    ):
        """Add a simple text box to a slide.

        Args:
            slide: PowerPoint slide object.
            text: Text content.
            left, top, width, height: Position and size.
            font_size: Font size.
            bold: Whether to make text bold.
            italic: Whether to make text italic.
            alignment: Text alignment.
            word_wrap: Whether to wrap text.

        Returns:
            Created textbox shape.
        """
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = word_wrap

        para = tf.paragraphs[0]
        para.text = text
        para.font.size = font_size or DEFAULT_BODY_FONT_SIZE
        para.font.bold = bold
        para.font.italic = italic
        para.alignment = alignment

        return shape

    def _fill_bullets(
        self,
        text_frame,
        items: List[dict],
        font_size: Optional[int] = None
    ) -> None:
        """Fill a text frame with bullet list content.

        This is the single method for rendering bullet lists, used by both
        placeholder-based slides and custom textbox-based slides.

        Supports inline markdown formatting in item text:
        **bold**, *italic*, ***bold italic***, ~~strikethrough~~,
        __underline__, `code`.

        Args:
            text_frame: PowerPoint text frame object (from placeholder or textbox).
            items: List of dicts with 'text' and 'indentation_level' keys.
            font_size: Optional font size for items.
        """
        if not items:
            return

        text_frame.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                para = text_frame.paragraphs[0]
            else:
                para = text_frame.add_paragraph()

            item_text = item.get("text", "")
            para.alignment = PP_ALIGN.LEFT
            para.level = max(0, int(item.get("indentation_level", 1)) - 1)

            # Apply inline markdown formatting if markers are present
            if has_inline_formatting(item_text):
                apply_inline_formatting(para, item_text, font_size=font_size)
            else:
                para.text = item_text
                if font_size:
                    para.font.size = font_size

    def _add_bullet_list(
        self,
        slide,
        items: List[dict],
        left: int,
        top: int,
        width: int,
        height: int,
        font_size: Optional[int] = None
    ):
        """Add a bullet list textbox to a slide.

        Args:
            slide: PowerPoint slide object.
            items: List of dicts with 'text' and 'indentation_level' keys.
            left, top, width, height: Position and size.
            font_size: Font size for items.

        Returns:
            Created textbox shape or None if no items.
        """
        if not items:
            return None

        shape = slide.shapes.add_textbox(left, top, width, height)
        self._fill_bullets(shape.text_frame, items, font_size)
        return shape

    # -------------------------------------------------------------------------
    # Table Helpers
    # -------------------------------------------------------------------------

    def _set_cell_fill(self, cell, color: RGBColor) -> None:
        """Set the background fill color of a table cell.

        Args:
            cell: Table cell object.
            color: RGBColor for the fill.
        """
        try:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()

            # Remove existing fill
            for child in list(tcPr):
                if child.tag.endswith('}solidFill'):
                    tcPr.remove(child)

            # Add new fill
            solidFill = parse_xml(
                f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<a:srgbClr val="{color}"/>'
                f'</a:solidFill>'
            )
            tcPr.append(solidFill)
        except Exception as e:
            logger.debug(f"Could not set cell fill color: {e}")

    def _create_styled_table(
        self,
        slide,
        table_data: List[List[str]],
        left: int,
        top: int,
        width: int,
        height: int,
        header_color: Optional[RGBColor] = None,
        alternate_rows: bool = True,
        column_alignments: Optional[List] = None
    ):
        """Create a styled table on a slide.

        Args:
            slide: PowerPoint slide object.
            table_data: List of rows (first row is header).
            left, top, width, height: Position and size.
            header_color: Header background color.
            alternate_rows: Whether to use alternating row colors.
            column_alignments: Optional list of PP_ALIGN values per column
                (extracted from markdown separator row).

        Returns:
            Created table shape.
        """
        num_rows = len(table_data)
        num_cols = max((len(row) for row in table_data), default=0)

        if num_rows == 0 or num_cols == 0:
            return None

        shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = shape.table

        header_color = header_color or TABLE_HEADER_FILL

        for row_idx, row_data in enumerate(table_data):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx >= num_cols:
                    continue

                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_text) if cell_text else ""

                # Apply column alignment
                if column_alignments and col_idx < len(column_alignments):
                    alignment = column_alignments[col_idx]
                    if alignment is not None:
                        cell.text_frame.paragraphs[0].alignment = alignment

                if row_idx == 0:  # Header row
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.color.rgb = TABLE_HEADER_TEXT
                    self._set_cell_fill(cell, header_color)
                elif alternate_rows and row_idx % 2 == 0:
                    self._set_cell_fill(cell, TABLE_ALT_ROW_FILL)

        return shape

    # -------------------------------------------------------------------------
    # Image Helpers
    # -------------------------------------------------------------------------

    def _add_image_from_url(
        self,
        slide,
        image_url: str,
        left: int,
        top: int,
        max_width: int,
        max_height: int,
        center_horizontal: bool = True,
        center_vertical: bool = False
    ) -> Optional[Any]:
        """Download and add an image from URL to a slide.

        Args:
            slide: PowerPoint slide object.
            image_url: URL of the image.
            left: Left position.
            top: Top position.
            max_width: Maximum width.
            max_height: Maximum height.
            center_horizontal: Whether to center horizontally.
            center_vertical: Whether to center vertically.

        Returns:
            Picture shape or None if failed.
        """
        if not image_url:
            return None

        try:
            image_data, _ = download_image(image_url)

            picture = slide.shapes.add_picture(
                image_data, left, top, width=max_width
            )

            # Scale to fit height if needed
            if picture.height > max_height:
                scale = max_height / picture.height
                picture.width = int(picture.width * scale)
                picture.height = max_height

            # Center if requested
            if center_horizontal:
                slide_width = self.presentation.slide_width
                picture.left = int((slide_width - picture.width) / 2)

            if center_vertical:
                picture.top = int(top + (max_height - picture.height) / 2)

            logger.debug(f"Added image from URL: {image_url}")
            return picture

        except (ImageDownloadError, ImageValidationError) as e:
            logger.error(f"Failed to download image: {e}")
            return None
        except Exception as e:
            logger.error(f"Failed to add image from URL '{image_url}': {e}", exc_info=True)
            return None

    def _add_image_placeholder(self, slide, message: str, left: int, top: int, width: int):
        """Add a placeholder text when image cannot be loaded.

        Args:
            slide: PowerPoint slide object.
            message: Error message to display.
            left, top, width: Position and width.
        """
        self._add_text_box(
            slide, f"[{message}]",
            left, top, width, Inches(1),
            italic=True, alignment=PP_ALIGN.CENTER
        )

