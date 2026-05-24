"""Tests for PowerPoint presentation creation.

These tests create actual .pptx files and save them to disk for manual inspection.
Output files are saved to tests/output/pptx/ directory.
"""

import sys
from pathlib import Path

# Add project root to path for imports
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

import pytest
from pptx_tools.slide_builder import PowerpointPresentation

# Output directory for test files
OUTPUT_DIR = Path(__file__).parent / "output" / "pptx"


@pytest.fixture(scope="module", autouse=True)
def setup_output_dir():
    """Create output directory if it doesn't exist."""
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    yield


def save_presentation(pres: PowerpointPresentation, filename: str) -> Path:
    """Save presentation to output directory and return path."""
    output_path = OUTPUT_DIR / filename
    buffer = pres.save()
    with open(output_path, "wb") as f:
        f.write(buffer.getvalue())
    print(f"Saved: {output_path}")
    return output_path


class TestBasicSlides:
    """Tests for basic slide types (title, section, content)."""

    def test_title_slide(self):
        """Test creating a presentation with a title slide."""
        slides = [
            {
                "slide_type": "title",
                "slide_title": "My Amazing Presentation",
                "subtitle": "Test Author"
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "01_title_slide.pptx")
        assert path.exists()
        assert path.stat().st_size > 0

    def test_section_slide(self):
        """Test creating a section slide."""
        slides = [
            {
                "slide_type": "title",
                "slide_title": "Presentation with Sections",
                "subtitle": "Tester"
            },
            {
                "slide_type": "section",
                "slide_title": "Part 1: Introduction"
            },
            {
                "slide_type": "section",
                "slide_title": "Part 2: Main Content"
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "02_section_slides.pptx")
        assert path.exists()

    def test_content_slide(self):
        """Test creating content slides with bullet points."""
        slides = [
            {
                "slide_type": "title",
                "slide_title": "Content Slides Demo",
                "subtitle": "Tester"
            },
            {
                "slide_type": "content",
                "slide_title": "Key Features",
                "slide_text": [
                    {"text": "First main point", "indentation_level": 1},
                    {"text": "Sub-point A", "indentation_level": 2},
                    {"text": "Sub-point B", "indentation_level": 2},
                    {"text": "Second main point", "indentation_level": 1},
                    {"text": "Third main point", "indentation_level": 1},
                    {"text": "Deep nested item", "indentation_level": 3}
                ]
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "03_content_slide.pptx")
        assert path.exists()


class TestTableSlides:
    """Tests for table slides with various configurations."""

    def test_basic_table(self):
        """Test creating a basic table slide."""
        slides = [
            {
                "slide_type": "table",
                "slide_title": "Sales Data Q1-Q4",
                "table_data": [
                    ["Product", "Q1", "Q2", "Q3", "Q4"],
                    ["Widget A", "100", "150", "200", "180"],
                    ["Widget B", "80", "90", "110", "130"],
                    ["Widget C", "200", "220", "250", "280"]
                ]
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "04_basic_table.pptx")
        assert path.exists()

    def test_styled_table(self):
        """Test table with custom header color and alternating rows."""
        slides = [
            {
                "slide_type": "table",
                "slide_title": "Styled Table",
                "table_data": [
                    ["Name", "Department", "Role", "Salary"],
                    ["John Doe", "Engineering", "Senior Dev", "$120,000"],
                    ["Jane Smith", "Marketing", "Manager", "$95,000"],
                    ["Bob Johnson", "Sales", "Rep", "$75,000"],
                    ["Alice Brown", "HR", "Director", "$110,000"],
                    ["Charlie Wilson", "Engineering", "Lead", "$140,000"]
                ],
                "header_color": "2E7D32",  # Green header
                "alternate_rows": True,
                "speaker_notes": "This table shows our team structure and compensation."
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "05_styled_table.pptx")
        assert path.exists()

    def test_table_no_alternating(self):
        """Test table without alternating row colors."""
        slides = [
            {
                "slide_type": "table",
                "slide_title": "Simple Table (No Zebra Stripes)",
                "table_data": [
                    ["Item", "Value"],
                    ["Alpha", "100"],
                    ["Beta", "200"],
                    ["Gamma", "300"]
                ],
                "alternate_rows": False
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "06_table_no_alternating.pptx")
        assert path.exists()

    def test_table_with_alignment(self):
        """Test table with markdown-style column alignment."""
        slides = [
            {
                "slide_type": "table",
                "slide_title": "Aligned Table",
                "table_data": [
                    ["Name", "Score", "Grade"],
                    [":---", ":---:", "---:"],  # left, center, right
                    ["Alice", "95", "A+"],
                    ["Bob", "82", "B"],
                ]
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "06b_aligned_table.pptx")
        assert path.exists()

        from pptx import Presentation as PptxReader
        from pptx.enum.text import PP_ALIGN
        prs = PptxReader(str(path))
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # Separator row should be removed: header + 2 data rows = 3
                assert len(table.rows) == 3
                # Column 1 should be centered
                assert table.cell(0, 1).text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER
                # Column 2 should be right-aligned
                assert table.cell(0, 2).text_frame.paragraphs[0].alignment == PP_ALIGN.RIGHT
                # Column 0 left (None = default)
                assert table.cell(0, 0).text_frame.paragraphs[0].alignment is None
                return
        assert False, "Table not found on slide"


class TestTwoColumnSlides:
    """Tests for two-column layout slides."""

    def test_two_column_basic(self):
        """Test basic two-column slide."""
        slides = [
            {
                "slide_type": "two_column",
                "slide_title": "Two Column Layout",
                "left_column": [
                    {"text": "Left point 1", "indentation_level": 1},
                    {"text": "Left point 2", "indentation_level": 1},
                    {"text": "Left sub-point", "indentation_level": 2}
                ],
                "right_column": [
                    {"text": "Right point 1", "indentation_level": 1},
                    {"text": "Right point 2", "indentation_level": 1}
                ]
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "07_two_column_basic.pptx")
        assert path.exists()

    def test_two_column_with_headings(self):
        """Test two-column slide with column headings."""
        slides = [
            {
                "slide_type": "two_column",
                "slide_title": "Pros and Cons Analysis",
                "left_heading": "✅ Advantages",
                "right_heading": "❌ Disadvantages",
                "left_column": [
                    {"text": "Faster development", "indentation_level": 1},
                    {"text": "Lower costs", "indentation_level": 1},
                    {"text": "Better scalability", "indentation_level": 1}
                ],
                "right_column": [
                    {"text": "Steeper learning curve", "indentation_level": 1},
                    {"text": "Limited documentation", "indentation_level": 1}
                ],
                "speaker_notes": "Emphasize the advantages outweigh disadvantages"
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "08_two_column_headings.pptx")
        assert path.exists()



class TestQuoteSlides:
    """Tests for quote slides."""

    def test_quote_with_author(self):
        """Test quote slide with author."""
        slides = [
            {
                "slide_type": "quote",
                "slide_title": "",
                "quote_text": "The only way to do great work is to love what you do.",
                "quote_author": "Steve Jobs",
                "speaker_notes": "Pause for effect after showing this quote"
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "10_quote_slide.pptx")
        assert path.exists()

    def test_quote_with_title(self):
        """Test quote slide with a title."""
        slides = [
            {
                "slide_type": "quote",
                "slide_title": "Inspirational Quote",
                "quote_text": "Innovation distinguishes between a leader and a follower.",
                "quote_author": "Steve Jobs"
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "11_quote_with_title.pptx")
        assert path.exists()


class TestChartSlides:
    """Tests for chart slides with various chart types."""

    def test_bar_chart(self):
        """Test bar chart slide."""
        slides = [
            {
                "slide_type": "chart",
                "slide_title": "Quarterly Sales (Bar Chart)",
                "chart_type": "bar",
                "chart_data": {
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": [
                        {"name": "2024", "values": [100, 150, 200, 180]},
                        {"name": "2025", "values": [120, 180, 220, 250]}
                    ]
                },
                "has_legend": True,
                "legend_position": "right"
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "12_bar_chart.pptx")
        assert path.exists()

    def test_column_chart(self):
        """Test column chart slide."""
        slides = [
            {
                "slide_type": "chart",
                "slide_title": "Revenue by Region",
                "chart_type": "column",
                "chart_data": {
                    "categories": ["North", "South", "East", "West"],
                    "series": [
                        {"name": "Revenue", "values": [450, 320, 280, 390]}
                    ]
                }
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "13_column_chart.pptx")
        assert path.exists()

    def test_line_chart(self):
        """Test line chart slide."""
        slides = [
            {
                "slide_type": "chart",
                "slide_title": "Monthly Trend",
                "chart_type": "line_markers",
                "chart_data": {
                    "categories": ["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
                    "series": [
                        {"name": "Users", "values": [1000, 1200, 1150, 1400, 1600, 1800]},
                        {"name": "Sessions", "values": [3000, 3500, 3200, 4200, 4800, 5500]}
                    ]
                },
                "legend_position": "bottom"
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "14_line_chart.pptx")
        assert path.exists()

    def test_pie_chart(self):
        """Test pie chart slide."""
        slides = [
            {
                "slide_type": "chart",
                "slide_title": "Market Share",
                "chart_type": "pie",
                "chart_data": {
                    "categories": ["Product A", "Product B", "Product C", "Others"],
                    "series": [
                        {"name": "Share", "values": [35, 28, 22, 15]}
                    ]
                },
                "speaker_notes": "Product A leads the market"
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "15_pie_chart.pptx")
        assert path.exists()

    def test_doughnut_chart(self):
        """Test doughnut chart slide."""
        slides = [
            {
                "slide_type": "chart",
                "slide_title": "Budget Allocation",
                "chart_type": "doughnut",
                "chart_data": {
                    "categories": ["R&D", "Marketing", "Operations", "HR"],
                    "series": [
                        {"name": "Budget", "values": [40, 25, 25, 10]}
                    ]
                }
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "16_doughnut_chart.pptx")
        assert path.exists()

    def test_stacked_bar_chart(self):
        """Test stacked bar chart slide."""
        slides = [
            {
                "slide_type": "chart",
                "slide_title": "Sales by Category (Stacked)",
                "chart_type": "bar_stacked",
                "chart_data": {
                    "categories": ["Region A", "Region B", "Region C"],
                    "series": [
                        {"name": "Electronics", "values": [100, 80, 120]},
                        {"name": "Clothing", "values": [60, 90, 70]},
                        {"name": "Food", "values": [40, 50, 45]}
                    ]
                }
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "17_stacked_bar_chart.pptx")
        assert path.exists()

    def test_area_chart(self):
        """Test area chart slide."""
        slides = [
            {
                "slide_type": "chart",
                "slide_title": "Growth Over Time",
                "chart_type": "area",
                "chart_data": {
                    "categories": ["2020", "2021", "2022", "2023", "2024"],
                    "series": [
                        {"name": "Revenue", "values": [100, 150, 180, 220, 300]}
                    ]
                }
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "18_area_chart.pptx")
        assert path.exists()


@pytest.mark.network
class TestImageSlides:
    """Tests for image slides with real images from picsum.photos.

    These tests require network access and may be flaky due to external service availability.
    Run with: pytest -m network
    Skip with: pytest -m "not network"
    """

    def test_image_slide_with_caption(self):
        """Test image slide with a real image and caption."""
        slides = [
            {
                "slide_type": "image",
                "slide_title": "Beautiful Landscape",
                "image_url": "https://picsum.photos/800/600",
                "image_caption": "Random image from Picsum Photos",
                "speaker_notes": "This is a real image downloaded from the internet"
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "19_image_with_caption.pptx")
        assert path.exists()
        assert path.stat().st_size > 40000  # Should be larger due to embedded image

    def test_image_slide_no_caption(self):
        """Test image slide without caption."""
        slides = [
            {
                "slide_type": "image",
                "slide_title": "Product Photo",
                "image_url": "https://picsum.photos/600/400"
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "20_image_no_caption.pptx")
        assert path.exists()

    def test_image_slide_portrait(self):
        """Test image slide with portrait orientation image."""
        slides = [
            {
                "slide_type": "image",
                "slide_title": "Portrait Image",
                "image_url": "https://picsum.photos/400/600",
                "image_caption": "Portrait orientation (400x600)"
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "21_image_portrait.pptx")
        assert path.exists()

    def test_image_slide_square(self):
        """Test image slide with square image."""
        slides = [
            {
                "slide_type": "image",
                "slide_title": "Square Image",
                "image_url": "https://picsum.photos/500/500",
                "image_caption": "Square format (500x500)"
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "22_image_square.pptx")
        assert path.exists()

    def test_image_slide_invalid_url(self):
        """Test image slide with invalid URL (shows placeholder)."""
        slides = [
            {
                "slide_type": "image",
                "slide_title": "Image Slide (Placeholder)",
                "image_url": "https://invalid-url-for-testing.example.com/image.png",
                "image_caption": "This should show a placeholder message"
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "23_image_placeholder.pptx")
        assert path.exists()

    def test_image_slide_no_url(self):
        """Test image slide without URL."""
        slides = [
            {
                "slide_type": "image",
                "slide_title": "Image Slide (No URL)"
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "24_image_no_url.pptx")
        assert path.exists()



class TestSpeakerNotes:
    """Tests for speaker notes functionality."""

    def test_speaker_notes_on_various_slides(self):
        """Test that speaker notes work on different slide types."""
        slides = [
            {
                "slide_type": "title",
                "slide_title": "Presentation with Speaker Notes",
                "subtitle": "Presenter"
            },
            {
                "slide_type": "content",
                "slide_title": "Content with Notes",
                "slide_text": [
                    {"text": "Main point", "indentation_level": 1}
                ],
                "speaker_notes": "Remember to emphasize this point strongly!"
            },
            {
                "slide_type": "table",
                "slide_title": "Data Table",
                "table_data": [
                    ["A", "B"],
                    ["1", "2"]
                ],
                "speaker_notes": "Explain each column carefully."
            },
            {
                "slide_type": "chart",
                "slide_title": "Chart with Notes",
                "chart_type": "pie",
                "chart_data": {
                    "categories": ["A", "B", "C"],
                    "series": [{"name": "Data", "values": [30, 40, 30]}]
                },
                "speaker_notes": "Point out the equal distribution between A and C."
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "28_speaker_notes.pptx")
        assert path.exists()


class TestFormats:
    """Tests for different presentation formats."""

    def test_format_16_9(self):
        """Test 16:9 widescreen format."""
        slides = [
            {
                "slide_type": "title",
                "slide_title": "Widescreen Presentation",
                "subtitle": "16:9 Format"
            },
            {
                "slide_type": "content",
                "slide_title": "Widescreen Content",
                "slide_text": [
                    {"text": "This is a 16:9 presentation", "indentation_level": 1}
                ]
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "29_format_16_9.pptx")
        assert path.exists()

    def test_format_4_3(self):
        """Test 4:3 standard format."""
        slides = [
            {
                "slide_type": "title",
                "slide_title": "Standard Presentation",
                "subtitle": "4:3 Format"
            },
            {
                "slide_type": "content",
                "slide_title": "Standard Content",
                "slide_text": [
                    {"text": "This is a 4:3 presentation", "indentation_level": 1}
                ]
            }
        ]
        pres = PowerpointPresentation(slides, "4:3")
        path = save_presentation(pres, "30_format_4_3.pptx")
        assert path.exists()


class TestCompletePresentation:
    """Test creating a complete presentation with all slide types."""

    def test_complete_presentation(self):
        """Create a comprehensive presentation demonstrating ALL available slide layouts.

        This test serves as a demonstration of all available slide types and layout options.
        Use the generated file (31_complete_presentation.pptx) as a visual reference.

        Layout mapping (PowerPoint default template):
        - Layout 0: Title Slide (title)
        - Layout 1: Title and Content (content, table, image, chart, quote)
        - Layout 2: Section Header (section)
        - Layout 3: Two Content (two_column without subheaders)
        - Layout 4: Comparison (two_column with subheaders)
        - Layout 5: Title Only
        - Layout 6: Blank
        """
        slides = [
            # =====================================================================
            # LAYOUT 0: Title Slide
            # =====================================================================
            {
                "slide_type": "title",
                "slide_title": "Complete PowerPoint Layout Demo",
                "subtitle": "MCP Office Documents - All Slide Types",
                "speaker_notes": "This presentation demonstrates all available slide layouts and options."
            },

            # =====================================================================
            # LAYOUT 2: Section Header
            # =====================================================================
            {
                "slide_type": "section",
                "slide_title": "Section 1: Basic Layouts"
            },

            # =====================================================================
            # LAYOUT 1: Title and Content - Bullet Points
            # =====================================================================
            {
                "slide_type": "content",
                "slide_title": "Content Slide with Bullet Points",
                "slide_text": [
                    {"text": "First level bullet point", "indentation_level": 1},
                    {"text": "Second level - more detail", "indentation_level": 2},
                    {"text": "Third level - even more detail", "indentation_level": 3},
                    {"text": "Back to first level", "indentation_level": 1},
                    {"text": "Another second level point", "indentation_level": 2}
                ],
                "speaker_notes": "Uses Layout 1 (Title and Content). Supports up to 3 indentation levels."
            },

            # =====================================================================
            # LAYOUT 2: Another Section Header
            # =====================================================================
            {
                "slide_type": "section",
                "slide_title": "Section 2: Data & Tables"
            },

            # =====================================================================
            # LAYOUT 1: Title and Content - Table
            # =====================================================================
            {
                "slide_type": "table",
                "slide_title": "Table Slide with Custom Styling",
                "table_data": [
                    ["Feature", "Basic", "Pro", "Enterprise"],
                    ["Users", "5", "25", "Unlimited"],
                    ["Storage", "10 GB", "100 GB", "1 TB"],
                    ["Support", "Email", "Priority", "24/7 Dedicated"],
                    ["Price", "$9/mo", "$29/mo", "Contact Us"]
                ],
                "header_color": "1565C0",
                "alternate_rows": True,
                "speaker_notes": "Uses Layout 1 (Title and Content). Table appears in content placeholder area."
            },

            # =====================================================================
            # LAYOUT 2: Another Section Header
            # =====================================================================
            {
                "slide_type": "section",
                "slide_title": "Section 3: Charts & Visualizations"
            },

            # =====================================================================
            # LAYOUT 1: Title and Content - Column Chart
            # =====================================================================
            {
                "slide_type": "chart",
                "slide_title": "Column Chart - Revenue Comparison",
                "chart_type": "column",
                "chart_data": {
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": [
                        {"name": "2023", "values": [200, 250, 300, 250]},
                        {"name": "2024", "values": [300, 350, 400, 450]}
                    ]
                },
                "has_legend": True,
                "legend_position": "right",
                "speaker_notes": "Uses Layout 1. Chart types: bar, column, line, pie, doughnut, stacked_bar, area"
            },

            # =====================================================================
            # LAYOUT 1: Title and Content - Pie Chart
            # =====================================================================
            {
                "slide_type": "chart",
                "slide_title": "Pie Chart - Market Share",
                "chart_type": "pie",
                "chart_data": {
                    "categories": ["Product A", "Product B", "Product C", "Other"],
                    "series": [{"name": "Market Share", "values": [35, 28, 22, 15]}]
                },
                "has_legend": True,
                "legend_position": "bottom",
                "speaker_notes": "Pie charts work best with a single series."
            },

            # =====================================================================
            # LAYOUT 1: Title and Content - Line Chart
            # =====================================================================
            {
                "slide_type": "chart",
                "slide_title": "Line Chart - Trend Analysis",
                "chart_type": "line",
                "chart_data": {
                    "categories": ["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
                    "series": [
                        {"name": "Website", "values": [1000, 1200, 1100, 1400, 1600, 1800]},
                        {"name": "Mobile App", "values": [500, 600, 800, 900, 1100, 1300]}
                    ]
                },
                "speaker_notes": "Line charts are great for showing trends over time."
            },

            # =====================================================================
            # LAYOUT 2: Another Section Header
            # =====================================================================
            {
                "slide_type": "section",
                "slide_title": "Section 4: Two Column Layouts"
            },

            # =====================================================================
            # LAYOUT 3: Two Content (without subheaders)
            # =====================================================================
            {
                "slide_type": "two_column",
                "slide_title": "Two Column Layout - No Subheaders",
                "left_column": [
                    {"text": "Left column content", "indentation_level": 1},
                    {"text": "More left content", "indentation_level": 1},
                    {"text": "Nested point", "indentation_level": 2}
                ],
                "right_column": [
                    {"text": "Right column content", "indentation_level": 1},
                    {"text": "More right content", "indentation_level": 1},
                    {"text": "Another nested point", "indentation_level": 2}
                ],
                "speaker_notes": "Uses Layout 3 (Two Content). No subheaders - just title and two content areas."
            },

            # =====================================================================
            # LAYOUT 4: Comparison (with subheaders)
            # =====================================================================
            {
                "slide_type": "two_column",
                "slide_title": "Comparison Layout - With Subheaders",
                "left_heading": "Before",
                "right_heading": "After",
                "left_column": [
                    {"text": "Manual processes", "indentation_level": 1},
                    {"text": "Time-consuming", "indentation_level": 2},
                    {"text": "Error-prone", "indentation_level": 2},
                    {"text": "Limited scalability", "indentation_level": 1}
                ],
                "right_column": [
                    {"text": "Automated workflows", "indentation_level": 1},
                    {"text": "Fast execution", "indentation_level": 2},
                    {"text": "Consistent results", "indentation_level": 2},
                    {"text": "Infinitely scalable", "indentation_level": 1}
                ],
                "speaker_notes": "Uses Layout 4 (Comparison). Has subheaders above each column for labeling."
            },

            # =====================================================================
            # LAYOUT 4: Another Comparison Example
            # =====================================================================
            {
                "slide_type": "two_column",
                "slide_title": "Pros and Cons Analysis",
                "left_heading": "✅ Advantages",
                "right_heading": "❌ Disadvantages",
                "left_column": [
                    {"text": "Easy to use", "indentation_level": 1},
                    {"text": "Cost effective", "indentation_level": 1},
                    {"text": "Flexible deployment", "indentation_level": 1}
                ],
                "right_column": [
                    {"text": "Learning curve", "indentation_level": 1},
                    {"text": "Initial setup time", "indentation_level": 1},
                    {"text": "Requires internet", "indentation_level": 1}
                ],
                "speaker_notes": "Comparison layout is perfect for pros/cons, before/after, or any side-by-side comparison."
            },

            # =====================================================================
            # LAYOUT 2: Another Section Header
            # =====================================================================
            {
                "slide_type": "section",
                "slide_title": "Section 5: Images & Quotes"
            },

            # =====================================================================
            # LAYOUT 1: Title and Content - Image (placeholder - no network needed)
            # =====================================================================
            {
                "slide_type": "image",
                "slide_title": "Image Slide with Caption",
                "image_url": "file:///nonexistent/sample.png",
                "image_caption": "Image placeholder (no network call in tests)",
                "speaker_notes": "Uses Layout 1. Images are downloaded, scaled to fit, and centered."
            },

            # =====================================================================
            # LAYOUT 1: Title and Content - Image without title
            # =====================================================================
            {
                "slide_type": "image",
                "image_url": "file:///nonexistent/large.png",
                "image_caption": "Full-bleed image slide (no title)",
                "speaker_notes": "Image slides can omit the title for a more impactful visual."
            },

            # =====================================================================
            # LAYOUT 1: Title and Content - Quote with title
            # =====================================================================
            {
                "slide_type": "quote",
                "slide_title": "Inspirational Quote",
                "quote_text": "The best way to predict the future is to create it.",
                "quote_author": "Peter Drucker",
                "speaker_notes": "Uses Layout 1. Quote slides can have an optional title."
            },

            # =====================================================================
            # LAYOUT 1: Title and Content - Quote without title
            # =====================================================================
            {
                "slide_type": "quote",
                "quote_text": "Innovation distinguishes between a leader and a follower.",
                "quote_author": "Steve Jobs",
                "speaker_notes": "Quote slides without title give more prominence to the quote itself."
            },

            # =====================================================================
            # LAYOUT 2: Final Section
            # =====================================================================
            {
                "slide_type": "section",
                "slide_title": "Conclusion"
            },

            # =====================================================================
            # LAYOUT 1: Summary Content Slide
            # =====================================================================
            {
                "slide_type": "content",
                "slide_title": "Available Slide Types Summary",
                "slide_text": [
                    {"text": "title - Opening slide (Layout 0: Title Slide)", "indentation_level": 1},
                    {"text": "section - Section dividers (Layout 2: Section Header)", "indentation_level": 1},
                    {"text": "content - Bullet points (Layout 1: Title and Content)", "indentation_level": 1},
                    {"text": "table - Data tables (Layout 1: Title and Content)", "indentation_level": 1},
                    {"text": "chart - Bar, column, line, pie, doughnut charts (Layout 1)", "indentation_level": 1},
                    {"text": "two_column - Side by side content (Layout 3 or 4)", "indentation_level": 1},
                    {"text": "Without subheaders → Two Content layout", "indentation_level": 2},
                    {"text": "With subheaders → Comparison layout", "indentation_level": 2},
                    {"text": "image - Images from URL (Layout 1: Title and Content)", "indentation_level": 1},
                    {"text": "quote - Quotations (Layout 1: Title and Content)", "indentation_level": 1}
                ],
                "speaker_notes": "All content-based slides use the Title and Content layout for consistent positioning."
            },

            # =====================================================================
            # LAYOUT 0: Closing Title Slide
            # =====================================================================
            {
                "slide_type": "title",
                "slide_title": "Thank You!",
                "subtitle": "Questions & Discussion"
            }
        ]

        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "31_complete_presentation.pptx")
        assert path.exists()
        print(f"\n✅ Complete presentation saved to: {path}")
        print(f"   File size: {path.stat().st_size / 1024:.1f} KB")
        print(f"   Total slides: {len(slides)}")
        print(f"\n   Slide types demonstrated:")
        print(f"   - title (Layout 0)")
        print(f"   - section (Layout 2)")
        print(f"   - content (Layout 1)")
        print(f"   - table (Layout 1)")
        print(f"   - chart - column, pie, line (Layout 1)")
        print(f"   - two_column without subheaders (Layout 3)")
        print(f"   - two_column with subheaders (Layout 4)")
        print(f"   - image (Layout 1)")
        print(f"   - quote (Layout 1)")


class TestEdgeCases:
    """Tests for edge cases and error handling."""

    def test_empty_content(self):
        """Test slide with empty content arrays."""
        slides = [
            {
                "slide_type": "two_column",
                "slide_title": "Empty Columns",
                "left_column": [],
                "right_column": []
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "32_empty_content.pptx")
        assert path.exists()

    def test_long_text(self):
        """Test handling of long text content."""
        long_text = "This is a very long piece of text that should wrap properly within the slide. " * 5
        slides = [
            {
                "slide_type": "content",
                "slide_title": "Long Text Handling",
                "slide_text": [
                    {"text": long_text, "indentation_level": 1}
                ]
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "33_long_text.pptx")
        assert path.exists()

    def test_special_characters(self):
        """Test handling of special characters."""
        slides = [
            {
                "slide_type": "content",
                "slide_title": "Special Characters: <>&\"'",
                "slide_text": [
                    {"text": "Arrows: → ← ↑ ↓", "indentation_level": 1},
                    {"text": "Math: ≤ ≥ ≠ ± × ÷", "indentation_level": 1},
                    {"text": "Currency: $ € £ ¥", "indentation_level": 1},
                    {"text": "Emoji: ✅ ❌ ⭐ 🚀", "indentation_level": 1}
                ]
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "34_special_characters.pptx")
        assert path.exists()

    def test_many_slides(self):
        """Test creating a presentation with many slides."""
        slides = [
            {"slide_type": "title", "slide_title": "Many Slides Test", "subtitle": "Tester"}
        ]
        for i in range(20):
            slides.append({
                "slide_type": "content",
                "slide_title": f"Slide {i + 1}",
                "slide_text": [
                    {"text": f"Content for slide {i + 1}", "indentation_level": 1}
                ]
            })
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "35_many_slides.pptx")
        assert path.exists()
        print(f"Created presentation with {len(slides)} slides")


class TestInlineFormatting:
    """Tests for inline markdown formatting in slide text."""

    def test_bold_and_italic(self):
        """Test bold and italic formatting in bullet points."""
        slides = [
            {
                "slide_type": "content",
                "slide_title": "Inline Formatting",
                "slide_text": [
                    {"text": "This is **bold** text", "indentation_level": 1},
                    {"text": "This is *italic* text", "indentation_level": 1},
                    {"text": "This is ***bold italic*** text", "indentation_level": 1},
                ]
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "40_inline_bold_italic.pptx")
        assert path.exists()

        # Verify runs were created with formatting
        from pptx import Presentation as PptxReader
        prs = PptxReader(str(path))
        slide = prs.slides[0]
        # Find the content placeholder
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if len(para.runs) > 1 and "bold" in para.text:
                        # Should have at least: "This is " + "bold" + " text"
                        bold_runs = [r for r in para.runs if r.font.bold]
                        assert len(bold_runs) > 0, "Expected bold run not found"
                        return
        # If we get here, we didn't find any formatted paragraphs
        assert False, "No formatted paragraphs found"

    def test_strikethrough_and_underline(self):
        """Test strikethrough and underline formatting."""
        slides = [
            {
                "slide_type": "content",
                "slide_title": "More Formatting",
                "slide_text": [
                    {"text": "Has ~~strikethrough~~ text", "indentation_level": 1},
                    {"text": "Has __underlined__ text", "indentation_level": 1},
                    {"text": "Has `code` text", "indentation_level": 1},
                ]
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "41_inline_strike_underline.pptx")
        assert path.exists()

        from pptx import Presentation as PptxReader
        prs = PptxReader(str(path))
        slide = prs.slides[0]
        found_underline = False
        found_code = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for r in para.runs:
                        if r.font.underline:
                            found_underline = True
                        if r.font.name == 'Courier New':
                            found_code = True
        assert found_underline, "Underline formatting not found"
        assert found_code, "Code (Courier New) formatting not found"

    def test_formatting_in_two_column(self):
        """Test inline formatting works in two-column slides."""
        slides = [
            {
                "slide_type": "two_column",
                "slide_title": "Formatted Columns",
                "left_heading": "Left",
                "right_heading": "Right",
                "left_column": [
                    {"text": "**Bold** point", "indentation_level": 1},
                ],
                "right_column": [
                    {"text": "*Italic* point", "indentation_level": 1},
                ],
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "42_inline_two_column.pptx")
        assert path.exists()

        from pptx import Presentation as PptxReader
        prs = PptxReader(str(path))
        slide = prs.slides[0]
        found_bold = False
        found_italic = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for r in para.runs:
                        if r.font.bold and r.text == "Bold":
                            found_bold = True
                        if r.font.italic and r.text == "Italic":
                            found_italic = True
        assert found_bold, "Bold formatting not found in left column"
        assert found_italic, "Italic formatting not found in right column"

    def test_plain_text_no_extra_runs(self):
        """Test that plain text without markers doesn't create multiple runs."""
        slides = [
            {
                "slide_type": "content",
                "slide_title": "Plain Text",
                "slide_text": [
                    {"text": "No formatting here", "indentation_level": 1},
                ]
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "43_plain_text.pptx")
        assert path.exists()

        from pptx import Presentation as PptxReader
        prs = PptxReader(str(path))
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text == "No formatting here":
                        # Plain text should use para.text directly (no runs via add_run)
                        assert len(para.runs) <= 1
                        return
        assert False, "Expected paragraph 'No formatting here' not found"

    def test_escaped_characters(self):
        """Test that escaped markdown characters render as literals."""
        slides = [
            {
                "slide_type": "content",
                "slide_title": "Escapes",
                "slide_text": [
                    {"text": "Use \\*asterisks\\* without formatting", "indentation_level": 1},
                ]
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "44_escaped_chars.pptx")
        assert path.exists()

        from pptx import Presentation as PptxReader
        prs = PptxReader(str(path))
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if "asterisks" in para.text:
                        # Escaped asterisks should render as literal *
                        assert "*asterisks*" in para.text
                        # Should NOT have italic formatting
                        italic_runs = [r for r in para.runs if r.font.italic]
                        assert len(italic_runs) == 0, "Escaped chars should not produce italic"
                        return
        assert False, "Paragraph with escaped characters not found"


class TestSubtitle:
    """Tests for subtitle field on title slides."""

    def test_subtitle_renders(self):
        """Test that subtitle text appears in the subtitle placeholder."""
        slides = [
            {
                "slide_type": "title",
                "slide_title": "Main Title",
                "subtitle": "This is the subtitle"
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "45_subtitle.pptx")
        assert path.exists()

        from pptx import Presentation as PptxReader
        prs = PptxReader(str(path))
        slide = prs.slides[0]
        found = False
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                assert ph.text == "This is the subtitle"
                found = True
        assert found, "Subtitle placeholder not found"

    def test_empty_subtitle(self):
        """Test that omitting subtitle doesn't cause errors."""
        slides = [
            {
                "slide_type": "title",
                "slide_title": "Title Only",
            }
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "46_no_subtitle.pptx")
        assert path.exists()


class TestFooterAndSlideNumbers:
    """Tests for footer text and slide number features."""

    def test_footer_text(self):
        """Test that footer text appears on slides."""
        slides = [
            {"slide_type": "title", "slide_title": "Footer Test"},
            {"slide_type": "content", "slide_title": "Slide 2",
             "slide_text": [{"text": "Content", "indentation_level": 1}]},
        ]
        pres = PowerpointPresentation(slides, "16:9", footer_text="Acme Corp")
        path = save_presentation(pres, "47_footer_text.pptx")
        assert path.exists()

        from pptx import Presentation as PptxReader
        prs = PptxReader(str(path))
        # Check footer on at least one slide
        found_footer = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    if shape.placeholder_format.idx == 11:
                        assert "Acme Corp" in shape.text
                        found_footer = True
        assert found_footer, "Footer placeholder not found on any slide"

    def test_slide_numbers(self):
        """Test that slide number placeholders are added."""
        slides = [
            {"slide_type": "title", "slide_title": "Numbers Test"},
            {"slide_type": "section", "slide_title": "Section"},
        ]
        pres = PowerpointPresentation(slides, "16:9", show_slide_numbers=True)
        path = save_presentation(pres, "48_slide_numbers.pptx")
        assert path.exists()

        from pptx import Presentation as PptxReader
        prs = PptxReader(str(path))
        found_sldnum = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    if shape.placeholder_format.idx == 12:
                        found_sldnum = True
        assert found_sldnum, "Slide number placeholder not found"

    def test_footer_special_characters(self):
        """Test that footer handles XML special characters safely."""
        slides = [
            {"slide_type": "content", "slide_title": "Special Footer",
             "slide_text": [{"text": "Test", "indentation_level": 1}]},
        ]
        # Footer with characters that would break XML if not escaped
        pres = PowerpointPresentation(
            slides, "16:9",
            footer_text='Acme & Co. <Confidential> "2025"'
        )
        path = save_presentation(pres, "49_footer_special_chars.pptx")
        assert path.exists()

        from pptx import Presentation as PptxReader
        prs = PptxReader(str(path))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    if shape.placeholder_format.idx == 11:
                        assert '&' in shape.text
                        assert '<' in shape.text
                        return
        assert False, "Footer with special chars not found"

    def test_no_footer_no_numbers(self):
        """Test that nothing is added when features are disabled."""
        slides = [
            {"slide_type": "content", "slide_title": "Clean Slide",
             "slide_text": [{"text": "Content", "indentation_level": 1}]},
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "50_no_footer_numbers.pptx")
        assert path.exists()

        from pptx import Presentation as PptxReader
        prs = PptxReader(str(path))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    idx = shape.placeholder_format.idx
                    assert idx not in (11, 12), f"Found unexpected placeholder idx={idx}"


class TestAuthorMetadata:
    """Tests for author document metadata."""

    def test_author_in_properties(self):
        """Test that author is stored in document core properties."""
        slides = [
            {"slide_type": "title", "slide_title": "Metadata Test", "subtitle": "A Subtitle"},
        ]
        pres = PowerpointPresentation(slides, "16:9", author="Custom Author")
        path = save_presentation(pres, "51_author_metadata.pptx")
        assert path.exists()

        from pptx import Presentation as PptxReader
        prs = PptxReader(str(path))
        assert prs.core_properties.author == "Custom Author"

    def test_no_author(self):
        """Test that omitting author preserves template default (not overwritten)."""
        slides = [
            {"slide_type": "title", "slide_title": "No Author"},
        ]
        pres = PowerpointPresentation(slides, "16:9")
        path = save_presentation(pres, "52_no_author.pptx")
        assert path.exists()

        from pptx import Presentation as PptxReader
        prs = PptxReader(str(path))
        # When author param is not provided, it should NOT be "Custom Author"
        # (it may be empty or carry the template's default metadata)
        assert prs.core_properties.author != "Custom Author"


if __name__ == "__main__":
    # Run tests with verbose output
    pytest.main([__file__, "-v", "-s"])







